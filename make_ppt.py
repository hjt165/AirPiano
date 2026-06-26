"""
Air Piano PPT Generator
Generates a 10-slide presentation using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Constants ---
TEAL = RGBColor(0x00, 0x89, 0x7B)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

SCREENSHOT_PATH = os.path.join(os.path.dirname(__file__), "测试3", "1.jpg")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "AirPiano_Presentation.pptx")
TEMP_PATH = os.path.join(os.path.dirname(__file__), "AirPiano_Presentation_new.pptx")


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text):
    """Add a teal title bar at the top of the slide."""
    # Teal background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT


def add_body_text(slide, lines, left=0.8, top=1.6, width=11.5, height=5.5, font_size=18):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line["text"]
        p.font.size = Pt(line.get("size", font_size))
        p.font.bold = line.get("bold", False)
        p.font.color.rgb = line.get("color", DARK_GRAY)
        p.font.name = line.get("font", "Calibri")
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line.get("space_after", 6))
        p.space_before = Pt(line.get("space_before", 0))

        if line.get("indent"):
            p.level = line["indent"]


def add_subtitle_line(slide, text, top=1.5):
    """Add a subtitle/tagline below the title bar."""
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(top), Inches(11.5), Inches(0.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT


def make_slide_1(prs):
    """Cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # Big teal block
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Air Piano"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "基于 MediaPipe 的手势演奏"
    p2.font.size = Pt(28)
    p2.font.color.rgb = WHITE
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER

    # Info lines below the teal block
    info_lines = [
        {"text": "课程名称：人工智能图像应用", "size": 20, "color": DARK_GRAY, "align": PP_ALIGN.CENTER, "space_after": 12},
        {"text": "日期：2026 年 6 月", "size": 18, "color": LIGHT_GRAY, "align": PP_ALIGN.CENTER},
    ]
    add_body_text(slide, info_lines, top=5.0, height=2.0)


def make_slide_2(prs):
    """Project overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "项目概述")

    add_subtitle_line(slide, "一句话描述：用手指在空中弹钢琴，无需任何物理乐器。")

    lines = [
        {"text": "核心功能", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  前置摄像头实时手部追踪", "size": 18, "space_after": 6},
        {"text": "  ●  手指敲击动作检测", "size": 18, "space_after": 6},
        {"text": "  ●  对应钢琴音符播放", "size": 18, "space_after": 6},
        {"text": "  ●  琴键可视化叠加层", "size": 18, "space_after": 16},
        {"text": "技术栈", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  Android  /  Kotlin  /  MediaPipe  /  CameraX", "size": 18, "space_after": 6},
    ]
    add_body_text(slide, lines, top=2.3)


def make_slide_3(prs):
    """Architecture / data flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "技术架构")

    lines = [
        {"text": "数据流", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  摄像头 (CameraX)  →  ImageProxy", "size": 17, "space_after": 4},
        {"text": "      ↓", "size": 17, "space_after": 4},
        {"text": "  MediaPipe Hand Landmarker  →  21 个手部关键点", "size": 17, "space_after": 4},
        {"text": "      ↓", "size": 17, "space_after": 4},
        {"text": "  敲击检测状态机  →  TAP_DOWN / HOLD / RELEASE", "size": 17, "space_after": 4},
        {"text": "      ↓", "size": 17, "space_after": 4},
        {"text": "  ├── 音符播放 (MediaPlayer)", "size": 17, "space_after": 4},
        {"text": "  └── 琴键 UI 更新 (OverlayView)", "size": 17, "space_after": 16},
        {"text": "关键组件", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  GestureRecognizerHelper — MediaPipe 推理封装", "size": 17, "space_after": 5},
        {"text": "  ●  AirPianoManager — 状态机 + 音频播放", "size": 17, "space_after": 5},
        {"text": "  ●  PianoKeyDrawer — 琴键绘制", "size": 17, "space_after": 5},
        {"text": "  ●  OverlayView — 手部关键点 + 琴键叠加", "size": 17, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.6, height=5.5)


def make_slide_4(prs):
    """Hand tracking."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "核心功能 1 — 手部追踪")

    lines = [
        {"text": "MediaPipe Hand Landmarker", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  21 个手部关键点（手腕、指关节、指尖）", "size": 18, "space_after": 5},
        {"text": "  ●  实时检测，支持单手/双手", "size": 18, "space_after": 5},
        {"text": "  ●  归一化坐标 (0-1)，适配不同分辨率", "size": 18, "space_after": 16},
        {"text": "关键点索引", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  0: 手腕      4: 拇指尖      8: 食指尖", "size": 18, "space_after": 5},
        {"text": "  12: 中指尖    16: 无名指尖    20: 小指尖", "size": 18, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.6)


def make_slide_5(prs):
    """Tap detection state machine."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "核心功能 2 — 敲击检测")

    lines = [
        {"text": "状态机设计", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  UNKNOWN  →  TAP_DOWN  →  HOLD  →  RELEASE", "size": 18, "space_after": 4},
        {"text": "     ↑                                      │", "size": 18, "space_after": 4},
        {"text": "     └──────────────────────────────────────┘", "size": 18, "space_after": 14},
        {"text": "检测原理", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  计算指尖 Y 轴偏移量相对于手腕的变化", "size": 18, "space_after": 5},
        {"text": "  ●  offsetDiff = tipOffset - palmOffset", "size": 18, "space_after": 5},
        {"text": "  ●  阈值判断状态转换", "size": 18, "space_after": 14},
        {"text": "为什么用相对偏移？", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  消除手部整体移动的干扰", "size": 18, "space_after": 5},
        {"text": "  ●  只检测手指的独立敲击动作", "size": 18, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.6)


def make_slide_6(prs):
    """Note playback."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "核心功能 3 — 音符播放")

    lines = [
        {"text": "手指-音符映射", "size": 22, "bold": True, "color": TEAL, "space_after": 10},
        {"text": "  拇指 → Do (C)     食指 → Re (D)     中指 → Mi (E)", "size": 18, "space_after": 5},
        {"text": "  无名指 → Fa (F)     小指 → Sol (G)", "size": 18, "space_after": 14},
        {"text": "实现方式", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  5 个 MediaPlayer 实例预加载", "size": 18, "space_after": 5},
        {"text": "  ●  TAP_DOWN 时 seekTo(0) + start()", "size": 18, "space_after": 5},
        {"text": "  ●  支持快速重复敲击", "size": 18, "space_after": 14},
        {"text": "音频规格", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  44100Hz 采样率，16-bit PCM WAV，指数衰减包络", "size": 18, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.6)


def make_slide_7(prs):
    """Piano key visualization + screenshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "核心功能 4 — 琴键可视化")

    lines = [
        {"text": "UI 设计", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  屏幕底部 18% 区域绘制 5 个等宽琴键", "size": 18, "space_after": 5},
        {"text": "  ●  空闲：白色 + 深灰边框 + 音名标签", "size": 18, "space_after": 5},
        {"text": "  ●  按下：独特颜色 + 白色文字 + 音符反馈", "size": 18, "space_after": 12},
        {"text": "配色方案", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  Do: 红色 (#FF6B6B)     Re: 橙色 (#FFA07A)     Mi: 金色 (#FFD700)", "size": 17, "space_after": 5},
        {"text": "  Fa: 绿色 (#98FB98)     Sol: 蓝色 (#87CEEB)", "size": 17, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.6, width=6.0)

    # Insert screenshot on the right side
    if os.path.exists(SCREENSHOT_PATH):
        slide.shapes.add_picture(
            SCREENSHOT_PATH,
            Inches(7.2), Inches(1.6), height=Inches(5.2)
        )


def make_slide_8(prs):
    """Challenges and solutions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "挑战与解决方案")

    lines = [
        {"text": "挑战 1：手指位置抖动", "size": 20, "bold": True, "color": TEAL, "space_after": 4},
        {"text": "  问题：摄像头噪声导致关键点坐标帧间跳动", "size": 17, "space_after": 3},
        {"text": "  方案：5 帧加权移动平均平滑（SMOOTHING_FACTOR=0.8）", "size": 17, "space_after": 14},
        {"text": "挑战 2：区分有意敲击与自然移动", "size": 20, "bold": True, "color": TEAL, "space_after": 4},
        {"text": "  问题：手部整体移动也会改变指尖坐标", "size": 17, "space_after": 3},
        {"text": "  方案：使用相对偏移（指尖 vs 手腕），而非绝对位置", "size": 17, "space_after": 14},
        {"text": "挑战 3：声音延迟", "size": 20, "bold": True, "color": TEAL, "space_after": 4},
        {"text": "  问题：MediaPlayer 首次播放有加载延迟", "size": 17, "space_after": 3},
        {"text": "  方案：预加载所有音符 + seekTo(0) 实现即时重播", "size": 17, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.6)


def make_slide_9(prs):
    """Demo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "演示")

    lines = [
        {"text": "演示流程", "size": 24, "bold": True, "color": TEAL, "space_after": 14},
        {"text": "  1.  打开应用，展示摄像头画面和手部追踪", "size": 20, "space_after": 12},
        {"text": "  2.  依次弯曲每根手指，展示单音符播放", "size": 20, "space_after": 12},
        {"text": "  3.  同时弯曲多根手指，展示和弦效果", "size": 20, "space_after": 12},
        {"text": "  4.  调整阈值设置，展示灵敏度变化", "size": 20, "space_after": 20},
        {"text": "准备：提前录制演示视频作为备份", "size": 18, "color": LIGHT_GRAY, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.8, height=5.0)


def make_slide_10(prs):
    """Summary & thanks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "团队分工与总结")

    lines = [
        {"text": "项目分工", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  需求分析与 UI 设计：[姓名]", "size": 18, "space_after": 5},
        {"text": "  核心算法开发：[姓名]", "size": 18, "space_after": 5},
        {"text": "  音频处理与集成：[姓名]", "size": 18, "space_after": 5},
        {"text": "  测试与文档：[姓名]", "size": 18, "space_after": 14},
        {"text": "项目收获", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  MediaPipe 手部追踪技术实践", "size": 18, "space_after": 5},
        {"text": "  ●  实时手势识别状态机设计", "size": 18, "space_after": 5},
        {"text": "  ●  Android CameraX 与自定义 View 开发", "size": 18, "space_after": 14},
        {"text": "未来改进方向", "size": 22, "bold": True, "color": TEAL, "space_after": 8},
        {"text": "  ●  支持更多音符（完整八度）", "size": 18, "space_after": 5},
        {"text": "  ●  双手演奏支持", "size": 18, "space_after": 5},
        {"text": "  ●  录制与回放功能", "size": 18, "space_after": 5},
        {"text": "  ●  自定义音色", "size": 18, "space_after": 5},
    ]
    add_body_text(slide, lines, top=1.6, height=5.5)

    # Thank you
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.3), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "谢谢！"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    make_slide_1(prs)
    make_slide_2(prs)
    make_slide_3(prs)
    make_slide_4(prs)
    make_slide_5(prs)
    make_slide_6(prs)
    make_slide_7(prs)
    make_slide_8(prs)
    make_slide_9(prs)
    make_slide_10(prs)

    prs.save(TEMP_PATH)
    # Try to replace original; if locked, keep the new file
    import shutil
    try:
        if os.path.exists(OUTPUT_PATH):
            os.remove(OUTPUT_PATH)
        shutil.move(TEMP_PATH, OUTPUT_PATH)
        print(f"PPT saved to: {OUTPUT_PATH}")
    except PermissionError:
        print(f"PPT saved to: {TEMP_PATH} (original file was locked)")


if __name__ == "__main__":
    main()
